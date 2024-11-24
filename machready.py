import streamlit as st
import sqlite3
import google.generativeai as genai
import pandas as pd
from PyPDF2 import PdfReader
import io
import os
from pptx import Presentation
import requests
import bcrypt
import pytz
import uuid
import base64
from datetime import datetime, timedelta
genai.configure(api_key=st.secrets["API_KEY"])

def login_page():
    st.subheader("ThisistheFUTURE")
    st.markdown("<h1 style='text-align: center; color: #ff5733;'>PedoMUS</h1>", unsafe_allow_html=True)
    st.markdown("""
    ### Welcome to PedoMUS
    **PedoMUS** is a comprehensive educational dashboard that provides an array of tools to support effective learning. The platform allows users to track and manage attendance with ease, interact with uploaded PPTs and PDFs through an intelligent chatbot that can summarize and explain content, and generate multiple-choice questions based on notes to aid in self-assessment. Additionally, PedoMUS offers access to simulation tools that facilitate a deeper understanding of complex concepts. Together, these features create an integrated, user-friendly environment to enrich the educational experience.
    """)
    # Auto-playing and looping video using custom HTML
    st.write("### Watch this while you log in!")
    video_html = """
    <video width="700" autoplay loop muted>
        <source src="https://www.w3schools.com/html/mov_bbb.mp4" type="video/mp4">
        Your browser does not support the video tag.
    </video>
    """
    st.markdown(video_html, unsafe_allow_html=True)

    user_id = st.text_input("User ID (Numerical Only)")
    password = st.text_input("Password", type="password")

    if st.button("Login"):
        if not user_id or not password:
            st.warning("Please enter both User ID and Password.")
            return

        if not user_id.isdigit():
            st.error("User ID must be numerical.")
            return

        user_type = authenticate(user_id, password)
        if user_type:
            st.session_state["logged_in"] = True
            st.session_state["user_id"] = user_id
            st.session_state["user_role"] = user_type
            st.success(f"Login successful! Welcome, {user_type.capitalize()}.")
            st.session_state["username"] = user_id
            st.success(f"Welcome to future of education, {user_id}!")
            st.session_state.login_time = datetime.now(pytz.timezone(get_user_timezone()))
            st.session_state.login_status = True
            return True
        else:
            st.error("Invalid User ID or Password.")
            
def get_llminfo():
    st.sidebar.header("Options", divider='rainbow')
    model = st.sidebar.radio("Choose LLM:", ("gemini-1.5-pro", "gemini-1.5-flash", "gemini-1.5-standard", "gemini-1.5-advanced"))
    temperature = st.sidebar.slider("Temperature:", 0.0, 2.0, 1.0, 0.25)
    top_p = st.sidebar.slider("Top P:", 0.0, 1.0, 0.94, 0.01)
    max_tokens = st.sidebar.slider("Maximum Tokens:", 100, 5000, 2000, 100)
    top_k = st.sidebar.slider("Top K:", 0, 100, 50, 1)
    return model, temperature, top_p, max_tokens, top_k


def save_teacher_attendance(present_students):
    """Save the attendance of the selected students to the database."""
    try:
        date_today = datetime.now().strftime("%Y-%m-%d")
        conn = sqlite3.connect("attendance.db")
        cursor = conn.cursor()
        
        # Create a table for attendance records if it doesn't exist
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS attendance_records (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                date TEXT NOT NULL,
                student_name TEXT NOT NULL
            )
        """)

        # Insert attendance records for present students
        for student in present_students:
            cursor.execute("INSERT INTO attendance_records (date, student_name) VALUES (?, ?)", (date_today, student))
        
        conn.commit()
        st.success(f"Attendance recorded for {len(present_students)} students on {date_today}.")
    except Exception as e:
        st.error(f"An error occurred while recording attendance: {e}")
    finally:
        conn.close()

def materials_dashboard():
    st.title("Learning Materials")
    os.makedirs("uploaded_materials", exist_ok=True)
    if st.session_state["user_role"] == "teacher":
        # Teacher can upload files
        uploaded_file = st.file_uploader("Upload Materials", type=['pdf', 'docx', 'pptx'])
        if uploaded_file:
            save_material(uploaded_file)
    materials = fetch_uploaded_materials()
    if materials:
        st.subheader("Available Materials")
        for idx, (filename, upload_time) in enumerate(materials):
            file_path = os.path.join("uploaded_materials", filename)
            with open(file_path, "rb") as f:
                st.download_button(f"Download {filename}", f, file_name=filename, key=f"download-{idx}")
    else:
        st.info("No materials have been uploaded yet.")


def mark_attendance_dashboard():
    role = st.session_state["user_role"]
    username = st.session_state["user_id"]

    st.title(f"{role.capitalize()} Dashboard - Attendance")
    if role == "teacher":
        # Teacher selects students
        students = [f"Student {i+1}" for i in range(1234)]
        present_students = st.multiselect("Select students who are present", students)
        if st.button("Submit Attendance"):
            if present_students:
                save_teacher_attendance(present_students)  # Now defined
            else:
                st.warning("Please select at least one student.")


def get_user_role(username):
    """Determine if the user is a teacher or student based on the username."""
    if username.isdigit():  # Ensure username is numerical
        if len(username) == 5:
            return "teacher"
        elif len(username) < 5:
            return "student"
    return None

# Initialize database for "present_today"
def init_present_today_db():
    conn = sqlite3.connect("attendance.db")
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS present_today (
            username TEXT PRIMARY KEY,
            timestamp TEXT NOT NULL
        )
    """)
    conn.commit()
    conn.close()

# Validate the code
def validate_code(input_code):
    try:
        conn = sqlite3.connect("codes.db")
        cursor = conn.cursor()

        # Check the code against the database
        cursor.execute("SELECT expiration_time FROM codes WHERE code = ?", (input_code,))
        result = cursor.fetchone()

        if result:
            expiration_time = datetime.strptime(result[0], "%Y-%m-%d %H:%M:%S")
            if datetime.now() <= expiration_time:
                return True  # Code is valid

        return False  # Code is invalid or expired
    except sqlite3.Error as e:
        st.error(f"Error validating the code: {e}")
        return False  # Return False in case of database errors
    finally:
        conn.close()

def questions_page():
    st.subheader("Questions Page")
    st.markdown("""The Questions page allows you to upload PDFs and PPTs, automatically extracting the information and generating interactive MCQ questions for practice. To your left is parameter control for the LLM you chose to use.""")
    model, temperature, top_p, max_tokens, top_k = get_llminfo()

    uploaded_file = st.file_uploader("Upload a PDF or PPT file", type=["pdf", "ppt", "pptx"])

    if uploaded_file is not None:
        text = ""
        if uploaded_file.type == "application/pdf":
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            presentation = Presentation(uploaded_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            text = text.strip()

        if st.button("Generate MCQs"):
            if text:
                generation_config = {
                    "temperature": temperature,
                    "top_p": top_p,
                    "max_output_tokens": max_tokens,
                    "top_k": top_k
                }
                model_instance = genai.GenerativeModel(model_name=model, generation_config=generation_config)

                prompt = (
                    f"Please generate exactly 5 multiple-choice questions based on the following text:\n\n"
                    f"{text}\n\n"
                    "Each question should have 4 options labeled A, B, C, and D. "
                    "At the end of each question, specify the correct answer in the format:\n"
                    "Correct Answer: [A/B/C/D]\n\n"
                    "For example:\n"
                    "1. What is the capital of France?\n"
                    "A) Paris\n"
                    "B) London\n"
                    "C) Rome\n"
                    "D) Berlin\n"
                    "Correct Answer: A\n\n"
                    "Now generate 5 questions following this format:"
                )
                
                response = model_instance.generate_content([prompt])
                mcqs_with_answers = response.text.strip().split('\n\n')

                st.session_state.mcqs = []
                st.session_state.correct_answers = []

                for mcq in mcqs_with_answers:
                    lines = mcq.split('\n')
                    if len(lines) < 2:
                        continue
                    
                    question_text = lines[0].strip()
                    options = [option.strip() for option in lines[1:5] if option.strip()]

                    correct_answer_line = lines[-1] if len(lines) > 5 else ""
                    correct_answer = correct_answer_line.split(":")[-1].strip().upper()
                    correct_answer = correct_answer.replace("**", "")

                    if not options or not correct_answer:
                        continue
                    
                    st.session_state.mcqs.append((question_text, options))
                    st.session_state.correct_answers.append(correct_answer)

                st.session_state.user_answers = [None] * len(st.session_state.mcqs)

        if 'mcqs' in st.session_state:
            st.subheader("Generated MCQs:")
            
            for i, (question_text, options) in enumerate(st.session_state.mcqs):
                selected_option = st.radio(
                    question_text, 
                    options, 
                    key=f"question_{i}", 
                    index=options.index(st.session_state.user_answers[i]) if st.session_state.user_answers[i] in options else 0
                )
                st.session_state.user_answers[i] = selected_option

            if st.button("Submit Answers"):
                correct_answers_count = 0

                for i, selected_option in enumerate(st.session_state.user_answers):
                    normalized_selected_option = selected_option[0].upper()
                    normalized_correct_answer = st.session_state.correct_answers[i].upper()

                    if normalized_selected_option == normalized_correct_answer:
                        correct_answers_count += 1
                
                total_questions = len(st.session_state.mcqs)
                st.success(f"You got {correct_answers_count} out of {total_questions} correct!")


def reading_material_page():
    st.subheader("Reading Material Interaction")
    st.markdown("""The Reading Material page allows you to upload various types of media, enabling you to chat with a chatbot about the content for better understanding and clarity also providing additional sources to read. To your left is parameter control for the LLM you chose to use.""")
    model, temperature, top_p, max_tokens, top_k = get_llminfo()

    typepdf = st.radio("Select the type of media to interact with:", ("PDF", "Images", "Videos", "PPT"), index=0)

    if typepdf == "PDF":
        st.write("You selected PDF. Upload your files below.")
        uploaded_files = st.file_uploader("Choose one or more PDFs", type='pdf', accept_multiple_files=True)
        if uploaded_files:
            text = ""
            for pdf in uploaded_files:
                pdf_reader = PdfReader(pdf)
                for page in pdf_reader.pages:
                    text += page.extract_text()

            generation_config = {
                "temperature": temperature,
                "top_p": top_p,
                "max_output_tokens": max_tokens,
                "top_k": top_k,
                "response_mime_type": "text/plain",
            }
            model_instance = genai.GenerativeModel(
                model_name=model,
                generation_config=generation_config,
            )
            st.write(model_instance.count_tokens(text))
            question = st.text_input("Enter your question and hit return.")
            if question:
                response = model_instance.generate_content([question, text])
                st.write(response.text)

    elif typepdf == "Images":
        st.write("You selected Images. Upload your image file below.")
        image_file = st.file_uploader("Upload your image file.", type=["jpg", "jpeg", "png"])
        if image_file:
            temp_file_path = image_file.name
            with open(temp_file_path, "wb") as f:
                f.write(image_file.getbuffer())

            st.write("Uploading image...")
            uploaded_image = genai.upload_file(path=temp_file_path)
            while uploaded_image.state.name == "PROCESSING":
                time.sleep(5)
                uploaded_image = genai.get_file(uploaded_image.name)

            if uploaded_image.state.name == "FAILED":
                st.error("Failed to process image.")
                return

            st.write("Image uploaded successfully. Enter your prompt below.")
            prompt2 = st.text_input("Enter your prompt for the image.")
            if prompt2:
                generation_config = {
                    "temperature": temperature,
                    "top_p": top_p,
                    "max_output_tokens": max_tokens,
                    "top_k": top_k,
                }
                response_image = model_instance.generate_content([prompt2, uploaded_image])
                st.write(response_image.text)

    elif typepdf == "Videos":
        st.write("You selected Videos. Upload your video file below.")
        video_file = st.file_uploader("Upload your video file.", type=["mp4", "mov", "avi"])
        if video_file:
            temp_file_path = video_file.name
            with open(temp_file_path, "wb") as f:
                f.write(video_file.getbuffer())

            st.write("Uploading video...")
            uploaded_video = genai.upload_file(path=temp_file_path)
            while uploaded_video.state.name == "PROCESSING":
                time.sleep(5)
                uploaded_video = genai.get_file(uploaded_video.name)

            if uploaded_video.state.name == "FAILED":
                st.error("Failed to process video.")
                return
            
            st.write("Video uploaded successfully. Enter your prompt below.")
            prompt3 = st.text_input("Enter your prompt for the video.")
            if prompt3:
                model_instance = genai.GenerativeModel(model_name=model)
                response = model_instance.generate_content([uploaded_video, prompt3])
                st.markdown(response.text)
                genai.delete_file(uploaded_video.name)

    elif typepdf == "PPT":
        st.write("You selected PPT. Upload your PowerPoint file below.")
        uploaded_ppt = st.file_uploader("Choose a PPT file", type='pptx')
        if uploaded_ppt:
            text = ""
            presentation = Presentation(uploaded_ppt)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            text = text.strip()
            st.write("Extracted text from PowerPoint:")
            st.write(text)

            generation_config = {
                "temperature": temperature,
                "top_p": top_p,
                "max_output_tokens": max_tokens,
                "top_k": top_k,
                "response_mime_type": "text/plain",
            }
            model_instance = genai.GenerativeModel(
                model_name=model,
                generation_config=generation_config,
            )
            st.write(model_instance.count_tokens(text))
            question = st.text_input("Enter your question about the PPT content and hit return.")
            if question:
                response = model_instance.generate_content([question, text])
                st.write(response.text)

def simulation_page():
    st.subheader("Simulation Page")
    st.write("Select a simulation to view:")
    st.markdown("""The Simulations page provides various simulation from Phet, which are a great tool to help build foundational knowledge.""")
    st.set_page_config(layout="wide")
    simulations = [
        ("Gene Expression Essentials", """
            <iframe src="https://phet.colorado.edu/sims/html/gene-expression-essentials/latest/gene-expression-essentials_en.html"
                width="650"
                height="500"
                allowfullscreen>
            </iframe>
        """),
        
        ("Photosynthesis", """
            <iframe src="https://phet.colorado.edu/sims/cheerpj/photoelectric/latest/photoelectric.html?simulation=photoelectric"
        width="100%"
        height="600"
        allowfullscreen>
        </iframe> """),
        
        ("Solarsystem", """
            <iframe src="https://phet.colorado.edu/sims/html/my-solar-system/latest/my-solar-system_en.html"
                 width="700"
                 height="600"
                 allowfullscreen>
              </iframe>
        """),
        ("Beer's Law Lab", """
            <iframe src="https://phet.colorado.edu/sims/html/beers-law-lab/latest/beers-law-lab_en.html"
                width="650"
                height="500"
                allowfullscreen>
            </iframe>
        """),
        ("Kepler's Laws", """
            <iframe src="https://phet.colorado.edu/sims/html/keplers-laws/latest/keplers-laws_en.html"
                width="650"
                height="500"
                allowfullscreen>
            </iframe>
        """),
        ("Hooke's Law", """
            <iframe src="https://phet.colorado.edu/sims/html/hookes-law/latest/hookes-law_en.html"
                width="650"
                height="500"
                allowfullscreen>
            </iframe>
        """),
        ("pH Scale Basics", """
            <iframe src="https://phet.colorado.edu/sims/html/ph-scale-basics/latest/ph-scale-basics_en.html"
                width="650"
                height="500"
                allowfullscreen>
            </iframe>
        """)
    ]
    st.markdown(iframe_code, unsafe_allow_html=True)

    simulation_names = [sim[0] for sim in simulations]
    selected_simulation_name = st.selectbox("Select Simulation", simulation_names)
    selected_simulation = next(sim for sim in simulations if sim[0] == selected_simulation_name)

    st.write(f"**{selected_simulation[0]}**")
    st.components.v1.html(selected_simulation[1], height=600)

# Function to get user's timezone based on IP address
def get_user_timezone():
    try:
        response = requests.get("https://ipinfo.io/json")
        if response.status_code == 200:
            data = response.json()
            return data.get("timezone", "UTC")
    except:
        pass
    return "UTC"

# Function to display flip clock based on the user's time zone
def display_flip_clock():
    user_timezone = get_user_timezone()
    now = datetime.now(pytz.timezone(user_timezone))

    flipclock_html = f"""
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/flipclock/0.7.8/flipclock.min.css">
    <div id="flip-clock" style="display: flex; justify-content: center;"></div>
    
    <script src="https://cdnjs.cloudflare.com/ajax/libs/jquery/3.6.0/jquery.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/flipclock/0.7.8/flipclock.min.js"></script>
    <script type="text/javascript">
        document.addEventListener("DOMContentLoaded", function() {{
            var clock = new FlipClock(document.getElementById('flip-clock'), {{
                clockFace: 'TwentyFourHourClock',
                showSeconds: true
            }});
        }});
    </script>
    """
    
    st.components.v1.html(flipclock_html, height=150)

# Function to add JavaScript for detecting tab changes
def detect_tab_switch():
    js_code = """
    <script>
        document.addEventListener("visibilitychange", function() {
            if (document.hidden) {
                alert("You are moving away from this page! Please stay on this tab.");
            }
        });
    </script>
    """
    st.components.v1.html(js_code)

def display_session_timer():
    if "login_time" in st.session_state:
        # Calculate the elapsed time since login in seconds
        elapsed_time = (datetime.now(pytz.timezone(get_user_timezone())) - st.session_state.login_time).total_seconds()
        session_timer_html = f"""
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/flipclock/0.7.8/flipclock.min.css">
        <div id="session-timer" style="display: flex; justify-content: center;"></div>
        
        <script src="https://cdnjs.cloudflare.com/ajax/libs/jquery/3.6.0/jquery.min.js"></script>
        <script src="https://cdnjs.cloudflare.com/ajax/libs/flipclock/0.7.8/flipclock.min.js"></script>
        <script type="text/javascript">
            document.addEventListener("DOMContentLoaded", function() {{
                var timer = new FlipClock(document.getElementById('session-timer'), {{
                    clockFace: 'MinuteCounter',
                    autoStart: true
                }});
                // Set the timer to start from the elapsed time
                timer.setTime({int(elapsed_time)});
                timer.start();
            }});
        </script>
        """
        
        st.components.v1.html(session_timer_html, height=150)

# Initialize database for codes
def init_code_db():
    conn = sqlite3.connect("codes.db")
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS codes (
            code TEXT PRIMARY KEY,
            expiration_time TEXT NOT NULL
        )
    """)
    conn.commit()
    conn.close()

# Generate a unique code
def generate_unique_code():
    """Generate a unique attendance code and store it in the database with an expiration time."""
    code = str(uuid.uuid4())[:8]  # Generate a unique code (8 characters)
    expiration_time = datetime.now() + timedelta(minutes=6)  # Code expires in 6 minutes

    try:
        conn = sqlite3.connect("codes.db")
        cursor = conn.cursor()

        # Create the codes table if it doesn't exist
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS codes (
                code TEXT PRIMARY KEY,
                expiration_time TEXT NOT NULL
            )
        """)
        conn.commit()

        # Store the code and its expiration time
        cursor.execute("INSERT INTO codes (code, expiration_time) VALUES (?, ?)", 
                       (code, expiration_time.strftime("%Y-%m-%d %H:%M:%S")))
        conn.commit()
        conn.close()
        return code, expiration_time
    except Exception as e:
        print(f"Error generating unique code: {e}")
        return None, None

# Authenticate user by checking credentials in the database
def authenticate(user_id, password):
    try:
        conn = sqlite3.connect('users.db')
        cursor = conn.cursor()
        cursor.execute("SELECT password_hash, user_type FROM users WHERE user_id = ?", (user_id,))
        user = cursor.fetchone()
        conn.close()

        if user:
            hashed_password, user_type = user
            if bcrypt.checkpw(password.encode('utf-8'), hashed_password):
                return user_type
        return None
    except Exception as e:
        st.error(f"Error during authentication: {e}")
        return None

def save_material(uploaded_file):
    # Ensure directory exists
    os.makedirs("uploaded_materials", exist_ok=True)
    conn = sqlite3.connect("shared_materials.db")
    cursor = conn.cursor()
    upload_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    try:
        # Save file to directory
        file_path = os.path.join("uploaded_materials", uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        # Save metadata to the database
        cursor.execute("INSERT INTO materials (filename, upload_time) VALUES (?, ?)", 
                       (uploaded_file.name, upload_time))
        conn.commit()
        st.success(f"File '{uploaded_file.name}' uploaded successfully!")
    except Exception as e:
        st.error(f"Failed to upload file: {e}")
    finally:
        conn.close()

# Professor Dashboard (Core App Functionality)
def professor_dashboard():
    st.title("Professor Dashboard")

    # Logout button in the sidebar
    if st.sidebar.button("Logout"):
        st.session_state.logged_in = False
        st.session_state.pop('login_time', None)
        st.success("You have been logged out.")

    # Attendance section
    st.subheader("Take Attendance")

    # Simulated list of students
    students = [f"Student {i+1}" for i in range(1234)]
    present_students = st.multiselect("Select students who are present", students, key="attendance_multiselect")

    if st.button("Submit Attendance"):
        if present_students:
            save_teacher_attendance(present_students)  # Ensure this function is defined
        else:
            st.warning("Please select at least one student.")
    
     # Option to generate a unique attendance code
    if st.button("Generate Code"):
        unique_code, expiration_time = generate_unique_code()
        if unique_code:
            st.success(f"Generated Code: {unique_code}")
            st.write(f"Code Expires At: {expiration_time.strftime('%H:%M:%S')}")
        else:
            st.error("Failed to generate a code.")

    st.subheader("Upload Learning Materials")
    uploaded_file = st.file_uploader("Choose a file to upload", type=['pdf', 'docx', 'pptx'])
    if uploaded_file is not None:
        save_material(uploaded_file)

    # Option to view attendance records
    st.subheader("View Attendance Records")
    if st.button("Show Attendance Records"):
        try:
            conn = sqlite3.connect("attendance.db")
            cursor = conn.cursor()

            # Fetch all attendance records
            cursor.execute("SELECT date, student_name FROM attendance_records ORDER BY date DESC")
            attendance_records = cursor.fetchall()

            # Fetch all present today records
            cursor.execute("SELECT username, timestamp FROM present_today ORDER BY timestamp DESC")
            present_today_records = cursor.fetchall()
            conn.close()

            # Convert attendance records to a DataFrame
            if attendance_records:
                attendance_df = pd.DataFrame(attendance_records, columns=["Date", "Student Name"])
                st.write("### Attendance Records")
                st.write(attendance_df)
            else:
                st.info("No attendance records found.")

            # Convert present today records to a DataFrame
            if present_today_records:
                present_today_df = pd.DataFrame(present_today_records, columns=["Username", "Timestamp"])
                st.write("### Present Today Records")
                st.write(present_today_df)
            else:
                st.info("No present today records found.")

        except Exception as e:
            st.error(f"An error occurred while fetching attendance records: {e}")

    # Sidebar statistics (optional)
    st.subheader("Dashboard Statistics")
    total_students = len(students)
    attendance_rate = 89  # Placeholder

    col1, col2 = st.columns(2)
    with col1:
        st.metric(label="Total Students", value=f"{total_students}")
    with col2:
        st.metric(label="Attendance Rate", value=f"{attendance_rate}%")

# Database setup and login system
def init_db():
    """Initialize the unified database for all users."""
    conn = sqlite3.connect('users.db')  # Use a single database for both teachers and students
    cursor = conn.cursor()
    
    # Drop the users table if it exists (to reset schema)
    cursor.execute("DROP TABLE IF EXISTS users")
    
    # Recreate the users table with the correct schema
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
            user_id TEXT PRIMARY KEY,
            password_hash TEXT NOT NULL,
            user_type TEXT NOT NULL  -- 'teacher' or 'student'
        )
    """)
    conn.commit()
    conn.close()
    print("Database initialized successfully!")


def add_user_to_db(username, password, role):
    try:
        conn = sqlite3.connect('users.db')
        c = conn.cursor()
        hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
        c.execute("INSERT INTO users (user_id, password_hash, user_type) VALUES (?, ?, ?)", (username, hashed_password, role))
        conn.commit()
        print(f"User '{username}' added successfully.")
    except sqlite3.IntegrityError:
        print(f"User '{username}' already exists.")
    except Exception as e:
        print(f"An error occurred: {e}")
    finally:
        conn.close()


# Shared materials database setup
def init_materials_db():
    os.makedirs("uploaded_materials", exist_ok=True)  # Ensure the directory exists
    conn = sqlite3.connect("shared_materials.db")
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS materials (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT NOT NULL,
            upload_time TEXT NOT NULL
        )
    """)
    conn.commit()
    conn.close()

# Fetch uploaded materials
def fetch_uploaded_materials():
    init_materials_db()  # Ensure the database is initialized
    conn = sqlite3.connect("shared_materials.db")
    cursor = conn.cursor()
    cursor.execute("SELECT filename, upload_time FROM materials ORDER BY upload_time DESC")
    materials = cursor.fetchall()
    conn.close()
    return materials

@st.cache_data
def load_default_timetable():
    # Sample timetable data
    data = {
        "Time": ["09-10 AM", "10-11 AM", "11-12 AM", "12-01 PM", "01-02 PM", "02-03 PM", "03-04 PM", "04-05 PM"],
        "Monday": ["Lecture / G:All C:PEV112 / R: 56-703 / S:BO301"] * 8,
        "Tuesday": ["Lecture / G:All C:PEV112 / R: 56-703 / S:BO301"] * 8,
        "Wednesday": ["Practical / G:1 C:PEV112 / R: 56-703 / S:BO301"] * 8,
        "Thursday": [""] * 8,
        "Friday": [""] * 8,
        "Saturday": [""] * 8
    }
    return pd.DataFrame(data)

def load_course_info():
    # Sample course information
    course_data = {
        "CourseCode": ["BTY396", "BTY416", "BTY441", "BTY463", "BTY464", "BTY496", "BTY499", "BTY651", "ICT202B", "PEA402", "PESS01", "PEV112"],
        "CourseType": ["CR", "CR", "EM", "CR", "CR", "CR", "CR", "PW", "CR", "OM", "PE", "OM"],
        "CourseName": ["BIOSEPARATION ENGINEERING", "BIOSEPARATION ENGINEERING LABORATORY", "PHARMACEUTICAL ENGINEERING", 
                       "BIOINFORMATICS AND COMPUTATIONAL BIOLOGY", "BIOINFORMATICS AND COMPUTATIONAL BIOLOGY LABORATORY", 
                       "METABOLIC ENGINEERING", "SEMINAR ON SUMMER TRAINING", "QUALITY CONTROL AND QUALITY ASSURANCE", 
                       "AI, ML AND EMERGING TECHNOLOGIES", "ANALYTICAL SKILLS -II", "MENTORING - VII", "VERBAL ABILITY"],
        "Credits": [3, 1, 3, 2, 1, 2, 3, 3, 2, 4, 0, 3],
        "Faculty": ["Dr. Ajay Kumar", "Dr. Ajay Kumar", "Dr. Shashank Garg", "Dr. Anish Kumar", 
                    "Dr. Anish Kumar", "Dr. Shashank Garg", "", "Dr. Aarti Bains", 
                    "Dr. Piyush Kumar Yadav", "Kamal Deep", "", "Jaskiranjit Kaur"]
    }
    return pd.DataFrame(course_data)


def mark_attendance(username):
    try:
        conn = sqlite3.connect("attendance.db")
        cursor = conn.cursor()
        
        # Get the current timestamp
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # Insert a new attendance record into the attendance_records table
        cursor.execute("INSERT INTO attendance_records (date, student_name) VALUES (?, ?)", (timestamp, username))
        
        conn.commit()
        st.success(f"Attendance marked successfully for {username} at {timestamp}!")
    except Exception as e:
        st.error(f"Failed to mark attendance: {e}")
    finally:
        conn.close()


def student_dashboard():
    st.sidebar.title("Navigation")
    page = st.sidebar.selectbox("Choose a page", ("Home", "Simulation", "Reading Material", "Questions", "Attendance"))
    
    

    # Logout button in the sidebar
    if st.sidebar.button("Logout"):
        st.session_state.logged_in = False
        st.session_state.pop('login_time', None)
        st.success("You have been logged out.")

    # Call the function to detect tab switches
    detect_tab_switch()
    
    if page == "Home":
        home()
    elif page == "Simulation":
        st.title("Simulation")
        display_flip_clock()
        display_session_timer()
        st.write("Implement simulations")
        simulation_page()
    elif page == "Reading Material":
        st.title("Reading Material")
        display_flip_clock()
        display_session_timer()
        st.write("Here are some flashcards/reading material to engage students.")
        reading_material_page()
    elif page == "Questions":
        st.title("Questions")
        display_flip_clock()
        display_session_timer()
        st.write("Welcome to the Engaging Page.")
        questions_page()
    elif page == "Attendance":
        st.title("YOUR Attendance")
        display_flip_clock()
        display_session_timer()
        st.markdown('''The Attendance page is a time-locked page preventing acess to other websites while in use, and automatically marks your attendance when in class.''')
        Attendance()

    
def Attendance ():   
    st.title("Student Dashboard")
    init_present_today_db()
   # Ensure the login_time is timezone-aware
    if "login_time" not in st.session_state:
        user_timezone = pytz.timezone(get_user_timezone())
        st.session_state["login_time"] = datetime.now(user_timezone)

    # Get the current time in the user's timezone
    current_time = datetime.now(pytz.timezone(get_user_timezone()))

    # Calculate elapsed time
    elapsed_time = (current_time - st.session_state["login_time"]).total_seconds()

    # Display waiting information if less than 3 minutes
    if elapsed_time < 180:  # 3 minutes
        st.info(f"Please wait for {int(180 - elapsed_time)} seconds before entering the code.")
        st.stop()  # Stop the execution of the dashboard

    # Show the code input field after 3 minutes
    input_code = st.text_input("Enter the code provided by your teacher:")
    if st.button("Submit Code"):
        if validate_code(input_code):
            username = st.session_state["user_id"]
            mark_attendance(username)  # This will now create a new attendance record
        else:
            st.error("Invalid code. Please try again.")           

def home():
    st.title("Academic Schedule and Course Information")
    st.write("Welcome to the Home Page.")

    timetable_df = load_default_timetable()
    tab1, tab2, tab3 = st.tabs(["Weekly Schedule", "Course Information", "Learning Materials"])

    with tab1:
        st.subheader("Weekly Class Schedule")
        st.dataframe(timetable_df)

    with tab2:
        st.subheader("Course Information")
        course_info_df = load_course_info()
        st.dataframe(course_info_df)

    with tab3:
        st.subheader("Learning Materials")
        materials = fetch_uploaded_materials()
        if materials:
            for idx, (filename, upload_time) in enumerate(materials):
                st.markdown(f"**{filename}** uploaded on {upload_time}")
                file_path = os.path.join("uploaded_materials", filename)
                with open(file_path, "rb") as f:
                    # Add a unique key using filename and index
                    st.download_button(f"Download {filename}", f, file_name=filename, key=f"download-{idx}")
        else:
            st.info("No materials have been uploaded yet.")

def main():

    if st.session_state.get("logged_in"):
        role = st.session_state["user_role"]
        if role == "teacher":
            professor_dashboard()  # Redirect to teacher dashboard
        elif role == "student":
            student_dashboard()  # Redirect to student dashboard
        else:
            st.error("Unknown role. Please contact the administrator.")
    else:
        login_page()


def app():
    if "logged_in" not in st.session_state:
        st.session_state.logged_in = False

    main()

if __name__ == "__main__":
    app()
